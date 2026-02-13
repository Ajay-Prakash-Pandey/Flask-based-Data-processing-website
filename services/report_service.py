import io
import pandas as pd
from typing import Optional, Tuple, Dict, Any
from datetime import datetime

def generate_csv_report(df: pd.DataFrame, filename: str) -> Tuple[Optional[bytes], Optional[str]]:
    """Generate CSV report"""
    try:
        csv_buffer = io.BytesIO()
        df.to_csv(csv_buffer, index=False)  # type: ignore
        csv_buffer.seek(0)
        
        export_filename = f"{filename.replace('.csv', '')}_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"
        
        return csv_buffer.getvalue(), export_filename
    except Exception:
        return None, None

def generate_json_report(df: pd.DataFrame, filename: str) -> Tuple[Optional[bytes], Optional[str]]:
    """Generate JSON report"""
    try:
        json_data = df.to_json(orient='records', indent=2)  # type: ignore
        json_buffer = io.BytesIO(json_data.encode())
        json_buffer.seek(0)
        
        export_filename = f"{filename.replace('.csv', '')}_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
        
        return json_buffer.getvalue(), export_filename
    except Exception:
        return None, None

def generate_xlsx_report(df: pd.DataFrame, filename: str, stats: Optional[Dict[str, Any]] = None) -> Tuple[Optional[bytes], Optional[str]]:
    """Generate Excel report with data and statistics"""
    try:
        xlsx_buffer = io.BytesIO()
        
        with pd.ExcelWriter(xlsx_buffer, engine='openpyxl') as writer:  # type: ignore
            # Data sheet
            df.to_excel(writer, sheet_name='Data', index=False)  # type: ignore
            
            # Statistics sheet
            if stats:
                stats_data: Dict[str, list[str]] = {
                    'Metric': [],
                    'Value': []
                }
                for key, value in stats.items():
                    stats_data['Metric'].append(str(key))
                    stats_data['Value'].append(str(value))
                
                stats_df = pd.DataFrame(stats_data)
                stats_df.to_excel(writer, sheet_name='Statistics', index=False)  # type: ignore
        
        xlsx_buffer.seek(0)
        export_filename = f"{filename.replace('.csv', '')}_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        
        return xlsx_buffer.getvalue(), export_filename
    except ImportError:
        raise ImportError("Excel export requires openpyxl. Install with: pip install openpyxl")
    except Exception:
        return None, None

def generate_pdf_report(df: pd.DataFrame, filename: str, stats: Optional[Dict[str, Any]] = None) -> Tuple[Optional[bytes], Optional[str]]:
    """Generate PDF report with data summary and statistics"""
    try:
        from reportlab.lib import colors  # type: ignore
        from reportlab.lib.pagesizes import letter, A4  # type: ignore
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak  # type: ignore
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle  # type: ignore
        from reportlab.lib.units import inch  # type: ignore
        
        pdf_buffer = io.BytesIO()
        doc = SimpleDocTemplate(pdf_buffer, pagesize=letter, topMargin=0.5*inch, bottomMargin=0.5*inch)
        elements = []
        styles = getSampleStyleSheet()
        
        # Title
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#333333'),
            spaceAfter=30,
            alignment=1
        )
        elements: list[Any] = []
        elements.append(Paragraph("Data Processing Report", title_style))  # type: ignore
        elements.append(Spacer(1, 0.3*inch))  # type: ignore
        
        # Report info
        info_style = styles['Normal']
        elements.append(Paragraph(f"<b>File:</b> {filename}", info_style))  # type: ignore
        elements.append(Paragraph(f"<b>Generated:</b> {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}", info_style))  # type: ignore
        elements.append(Paragraph(f"<b>Rows:</b> {len(df)} | <b>Columns:</b> {len(df.columns)}", info_style))  # type: ignore
        elements.append(Spacer(1, 0.2*inch))  # type: ignore
        
        # Statistics section
        if stats:
            elements.append(Paragraph("<b>Data Summary</b>", styles['Heading2']))  # type: ignore
            stats_data = [['Metric', 'Value']]
            for key, value in list(stats.items())[:10]:
                stats_data.append([str(key), str(value)])
            
            stats_table = Table(stats_data, colWidths=[2.5*inch, 2*inch])
            stats_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#667eea')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 12),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('FONTSIZE', (0, 1), (-1, -1), 10),
            ]))
            elements.append(stats_table)  # type: ignore
            elements.append(Spacer(1, 0.2*inch))  # type: ignore
        
        # Data preview
        elements.append(PageBreak())  # type: ignore
        elements.append(Paragraph("<b>Data Preview (First 20 rows)</b>", styles['Heading2']))  # type: ignore
        
        preview_df = df.head(20)
        data_rows = [list(preview_df.columns)]
        for _, row in preview_df.iterrows():
            data_rows.append([str(x)[:30] for x in row])
        
        data_table = Table(data_rows, colWidths=[4.5*inch / len(preview_df.columns)] * len(preview_df.columns))
        data_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#667eea')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 10),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 10),
            ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
            ('FONTSIZE', (0, 1), (-1, -1), 8),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f0f0f0')]),
        ]))
        elements.append(data_table)  # type: ignore
        
        doc.build(elements)  # type: ignore
        pdf_buffer.seek(0)
        
        export_filename = f"{filename.replace('.csv', '')}_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        
        return pdf_buffer.getvalue(), export_filename
    except ImportError:
        raise ImportError("PDF export requires reportlab. Install with: pip install reportlab")
    except Exception:
        return None, None

def generate_ppt_report(df: pd.DataFrame, filename: str, stats: Optional[Dict[str, Any]] = None) -> Tuple[Optional[bytes], Optional[str]]:
    """Generate PowerPoint report"""
    try:
        from pptx import Presentation  # type: ignore
        from pptx.util import Inches, Pt  # type: ignore
        from pptx.enum.text import PP_ALIGN  # type: ignore
        from pptx.dml.color import RGBColor  # type: ignore
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Title
        slide1 = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide1.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(102, 126, 234)
        
        title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = "Data Processing Report"
        title_para.font.size = Pt(54)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER
        
        subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
        subtitle_para.font.size = Pt(20)
        subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
        subtitle_para.alignment = PP_ALIGN.CENTER
        
        # Slide 2: Data Summary
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide2.shapes.title
        if title is not None:  # type: ignore
            title.text = "Data Summary"  # type: ignore
            title.text_frame.paragraphs[0].font.size = Pt(44)  # type: ignore
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)  # type: ignore
        
        content_box = slide2.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        info_text = f"""File: {filename}
Total Rows: {len(df)}
Total Columns: {len(df.columns)}
Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}

Columns: {', '.join(df.columns[:5])}{'...' if len(df.columns) > 5 else ''}"""
        
        p = content_frame.paragraphs[0]
        p.text = info_text.strip()
        p.font.size = Pt(18)
        p.level = 0
        
        # Slide 3: Statistics
        if stats:
            slide3 = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide3.shapes.title
            if title is not None:  # type: ignore
                title.text = "Data Statistics"  # type: ignore
                title.text_frame.paragraphs[0].font.size = Pt(44)  # type: ignore
                title.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)  # type: ignore
            
            content_box = slide3.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5.5))
            content_frame = content_box.text_frame
            content_frame.word_wrap = True
            
            stats_text = ""
            for _, (key, value) in enumerate(list(stats.items())[:8]):
                stats_text += f"{key}: {value}\n"
            
            p = content_frame.paragraphs[0]
            p.text = stats_text.strip()
            p.font.size = Pt(16)
        
        # Slide 4: Data Preview
        slide4 = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide4.shapes.title
        if title is not None:  # type: ignore
            title.text = "Data Preview"  # type: ignore
            title.text_frame.paragraphs[0].font.size = Pt(44)  # type: ignore
            title.text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 126, 234)  # type: ignore
        
        preview_df = df.head(8)
        rows, cols = len(preview_df) + 1, len(preview_df.columns)
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)
        
        table_shape = slide4.shapes.add_table(rows, cols, left, top, width, height).table
        
        for col_idx, col_name in enumerate(preview_df.columns):
            cell = table_shape.cell(0, col_idx)
            cell.text = str(col_name)[:15]
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(102, 126, 234)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        for row_idx, (_, row) in enumerate(preview_df.iterrows(), 1):
            for col_idx, value in enumerate(row):
                cell = table_shape.cell(row_idx, col_idx)
                cell.text = str(value)[:20]
                cell.text_frame.paragraphs[0].font.size = Pt(9)
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
        
        ppt_buffer = io.BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        
        export_filename = f"{filename.replace('.csv', '')}_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        
        return ppt_buffer.getvalue(), export_filename
    except ImportError:
        raise ImportError("PowerPoint export requires python-pptx. Install with: pip install python-pptx")
    except Exception:
        return None, None

def generate_docx_report(df: pd.DataFrame, filename: str, stats: Optional[Dict[str, Any]] = None) -> Tuple[Optional[bytes], Optional[str]]:
    """Generate Word document report"""
    try:
        from docx import Document  # type: ignore
        from docx.shared import Inches, Pt, RGBColor  # type: ignore
        from docx.enum.text import WD_ALIGN_PARAGRAPH  # type: ignore
        
        doc = Document()
        
        # Title
        title = doc.add_heading('Data Processing Report', 0)
        title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        
        # Summary section
        doc.add_heading('Data Summary', level=1)
        
        summary_para = doc.add_paragraph()
        summary_para.add_run(f"File: ").bold = True
        summary_para.add_run(f"{filename}\n")
        summary_para.add_run(f"Generated: ").bold = True
        summary_para.add_run(f"{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
        summary_para.add_run(f"Total Rows: ").bold = True
        summary_para.add_run(f"{len(df)}\n")
        summary_para.add_run(f"Total Columns: ").bold = True
        summary_para.add_run(f"{len(df.columns)}\n")
        
        # Columns list
        doc.add_heading('Columns', level=2)
        columns_list = doc.add_paragraph(style='List Bullet')
        for col in df.columns:
            columns_list.add_run(f"{col}\n").italic = False
        
        # Statistics section
        if stats:
            doc.add_heading('Data Statistics', level=1)
            for key, value in list(stats.items())[:15]:
                p = doc.add_paragraph()
                p.add_run(f"{key}: ").bold = True
                p.add_run(f"{value}")
        
        # Data preview table
        doc.add_heading('Data Preview (First 10 rows)', level=1)
        
        preview_df = df.head(10)
        table = doc.add_table(rows=len(preview_df) + 1, cols=len(preview_df.columns))
        table.style = 'Light Grid Accent 1'
        
        header_cells = table.rows[0].cells
        for idx, col_name in enumerate(preview_df.columns):
            header_cells[idx].text = str(col_name)
            for paragraph in header_cells[idx].paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
        
        for row_idx, (_, row) in enumerate(preview_df.iterrows(), 1):
            row_cells = table.rows[row_idx].cells
            for col_idx, value in enumerate(row):
                row_cells[col_idx].text = str(value)[:50]
        
        # Footer
        doc.add_paragraph()
        footer = doc.add_paragraph("Generated using Flask Data Processing Application")
        footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
        footer_runs = footer.runs[0]
        footer_runs.font.size = Pt(10)
        footer_runs.font.italic = True
        footer_runs.font.color.rgb = RGBColor(128, 128, 128)
        
        docx_buffer = io.BytesIO()
        doc.save(docx_buffer)
        docx_buffer.seek(0)
        
        export_filename = f"{filename.replace('.csv', '')}_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
        
        return docx_buffer.getvalue(), export_filename
    except ImportError:
        raise ImportError("Word export requires python-docx. Install with: pip install python-docx")
    except Exception:
        return None, None
